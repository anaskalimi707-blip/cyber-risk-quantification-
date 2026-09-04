"""
CyberOptix Enterprise - Executive Pitch Deck Generator
Creates a 16:9 widescreen PowerPoint presentation (.pptx) with an executive Obsidian dark aesthetic.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck(output_filename="CyberOptix_Enterprise_Pitch_Deck.pptx"):
    prs = Presentation()
    # 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    c_bg = RGBColor(8, 13, 26)         # Deep Obsidian #080D1A
    c_card = RGBColor(15, 23, 42)      # Slate Card #0F172A
    c_card_border = RGBColor(30, 41, 59) # Slate Border #1E293B
    c_cyan = RGBColor(56, 189, 248)    # Sky Cyan #38BDF8
    c_teal = RGBColor(45, 212, 191)    # Mint Teal #2DD4BF
    c_amber = RGBColor(251, 191, 36)   # Gold Amber #FBBF24
    c_crimson = RGBColor(251, 113, 133)# Ruby Crimson #FB7185
    c_white = RGBColor(248, 250, 252)  # White #F8FAFC
    c_muted = RGBColor(148, 163, 184)  # Muted Slate #94A3B8

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = c_bg
        bg.line.fill.background()
        return bg

    def add_card(slide, left, top, width, height, bg_col=c_card, border_col=c_card_border):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_col
        card.line.color.rgb = border_col
        card.line.width = Pt(1.5)
        return card

    def add_header(slide, eyebrow, title, subtitle=None):
        txbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = eyebrow.upper()
        p1.font.size = Pt(10.5)
        p1.font.bold = True
        p1.font.color.rgb = c_cyan
        p1.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = c_white

        if subtitle:
            p3 = tf.add_paragraph()
            p3.text = subtitle
            p3.font.size = Pt(12)
            p3.font.color.rgb = c_muted
            p3.space_before = Pt(4)

    # -------------------------------------------------------------------------
    # SLIDE 1: Title & Executive Summary
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1)

    # Main Hero Title Box
    h_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(3.2))
    htf = h_box.text_frame
    htf.word_wrap = True
    
    p = htf.paragraphs[0]
    p.text = "CYBEROPTIX ENTERPRISE • CRIM-X APEX"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(10)

    p = htf.add_paragraph()
    p.text = "Know Your Cyber Risk in Money.\nInvest Where Every Rupee Reduces Risk."
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.space_after = Pt(14)

    p = htf.add_paragraph()
    p.text = "Continuous Financial Cyber-Risk Quantification (CRQ), Causal DML & Multi-Objective Capital Optimization"
    p.font.size = Pt(15)
    p.font.color.rgb = c_muted

    # 3 Stat Cards on Cover
    stats = [
        ("₹8.60 Cr -> ₹2.40 Cr", "Expected Annual Loss Reduction", c_crimson),
        ("2.8x Capital ROSI", "Return on Security Investment", c_teal),
        ("SEBI & NIST Aligned", "FAIR + CSCRF Grounded", c_cyan)
    ]
    for i, (val, lbl, col) in enumerate(stats):
        cx = Inches(1.0 + i * 3.85)
        add_card(s1, cx, Inches(5.0), Inches(3.6), Inches(1.5))
        s_box = s1.shapes.add_textbox(cx + Inches(0.2), Inches(5.2), Inches(3.2), Inches(1.1))
        stf = s_box.text_frame
        stf.word_wrap = True
        p1 = stf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = stf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(11)
        p2.font.color.rgb = c_muted
        p2.space_before = Pt(4)

    # -------------------------------------------------------------------------
    # SLIDE 2: The Problem - The Trillion-Dollar Translation Gap
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_bg(s2)
    add_header(s2, "Market Disconnect", "The Trillion-Dollar Translation Gap in Enterprise Security", "Why $200B in cybersecurity spending still fails to answer the Board's core fiduciary questions.")

    col_data = [
        ("The CISO's Dilemma", "Reports in technical CVEs & subjective red/yellow/green heatmaps. Unable to express cybersecurity in capital or business terms to the Board.", c_crimson),
        ("The CFO's Dilemma", "Approves spend with zero quantifiable ROI metrics. Security is perceived as an arbitrary insurance tax with no defensible proof of risk reduction.", c_amber),
        ("The Regulatory Crisis", "SEBI CSCRF, DPDP Act, and SEC mandate quantified board-level oversight. Annual manual questionnaires fail to catch continuous exposure drift.", c_cyan)
    ]
    for i, (title, desc, col) in enumerate(col_data):
        cx = Inches(0.8 + i * 3.95)
        add_card(s2, cx, Inches(2.2), Inches(3.75), Inches(4.5))
        c_box = s2.shapes.add_textbox(cx + Inches(0.3), Inches(2.5), Inches(3.15), Inches(3.8))
        ctf = c_box.text_frame
        ctf.word_wrap = True
        p1 = ctf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.space_after = Pt(12)
        p2 = ctf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = c_white
        p2.space_after = Pt(16)
        p3 = ctf.add_paragraph()
        p3.text = "Result: Misallocated capital, board blindness, and unchecked exposure."
        p3.font.size = Pt(11)
        p3.font.color.rgb = c_muted

    # -------------------------------------------------------------------------
    # SLIDE 3: The Solution - CyberOptix Enterprise Platform
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_bg(s3)
    add_header(s3, "Product Breakthrough", "The CyberOptix Enterprise Solution", "A single unified platform converting continuous telemetry into defensible financial risk & capital plans.")

    sol_cards = [
        ("Continuous FAIR Quantification", "Ingests telemetry from Okta, Qualys, CrowdStrike, and AWS. Runs 50,000 Monte Carlo iterations in <0.25s to output EAL & VaR in INR Lakh/Crore.", c_cyan),
        ("Causal AI (CRIM-X DML)", "Uses Robinson's Double Machine Learning to eliminate confounding bias and compute isolated causal risk reduction (ΔEAL) across 4,200 peer trials.", c_teal),
        ("5D Pareto Capital Optimizer", "Mixed-integer programming knapsack solver optimizing across Cost, Risk, Deployment SLA, Compliance Gain, and Disruption Index.", c_amber),
        ("Grounded AI Copilot", "Permission-checked natural language queries backed by an 11-tool API contract, cryptographic evidence citations, and strict human sign-off boundaries.", c_crimson)
    ]
    for i, (title, desc, col) in enumerate(sol_cards):
        col_idx = i % 2
        row_idx = i // 2
        cx = Inches(0.8 + col_idx * 5.95)
        cy = Inches(2.2 + row_idx * 2.4)
        add_card(s3, cx, cy, Inches(5.75), Inches(2.2))
        s_box = s3.shapes.add_textbox(cx + Inches(0.3), cy + Inches(0.25), Inches(5.15), Inches(1.7))
        stf = s_box.text_frame
        stf.word_wrap = True
        p1 = stf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.space_after = Pt(6)
        p2 = stf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = c_white

    # -------------------------------------------------------------------------
    # SLIDE 4: Core Moat - CRIM-X Causal DML vs Naive Correlation
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_bg(s4)
    add_header(s4, "Proprietary Moat", "CRIM-X Causal DML vs. Naive Correlation", "Why observational correlation in legacy tools misleads CFOs, and how Causal DML solves it.")

    # Left: Naive Correlation
    add_card(s4, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    lbox = s4.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(4.0))
    ltf = lbox.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "Naive Observational Correlation (Legacy)"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = c_crimson
    p.space_after = Pt(10)
    bullets_l = [
        "Claims: 'Deploying EDR reduced ransomware risk by 42%'.",
        "Confounded by simultaneous IT initiatives (cloud migration, patch SLAs, hiring).",
        "Double-counts overlapping control synergies.",
        "Overestimates risk reduction by up to 2.4x.",
        "Leads to bloated budgets on placebo controls."
    ]
    for b in bullets_l:
        p = ltf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = c_muted
        p.space_after = Pt(6)

    # Right: CRIM-X Causal DML
    add_card(s4, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), border_col=c_teal)
    rbox = s4.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(4.0))
    rtf = rbox.text_frame
    rtf.word_wrap = True
    p = rtf.paragraphs[0]
    p.text = "CRIM-X Causal Double Machine Learning (CyberOptix)"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = c_teal
    p.space_after = Pt(10)
    bullets_r = [
        "Robinson's DML separates confounding using orthogonalized Random Forest residuals.",
        "Cross-fitting (5-fold) eliminates regularization and overfitting bias.",
        "Identifies true isolated causal effect theta (ΔEAL in INR).",
        "Surfaces causal identification strategy (natural experiment, instrumental variable).",
        "Delivers mathematically defensible numbers to CFOs."
    ]
    for b in bullets_r:
        p = rtf.add_paragraph()
        p.text = "✓ " + b
        p.font.size = Pt(12)
        p.font.color.rgb = c_white
        p.space_after = Pt(6)

    # -------------------------------------------------------------------------
    # SLIDE 5: 5D Pareto Frontier Multi-Objective Optimization
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_bg(s5)
    add_header(s5, "Capital Allocation", "5D Multi-Objective Pareto Frontier Optimizer", "Moving beyond simple cost-benefit to balance enterprise deployment realities.")

    portfolios = [
        ("Balanced Frontier", "₹60 Lakh", "₹2.10 Cr", "2.8x", "Optimal balance between spend, speed, and risk reduction.", c_teal),
        ("Rapid Sprint", "₹45 Lakh", "₹1.65 Cr", "3.1x", "Fully deployable in under 21 days for urgent regulatory deadlines.", c_cyan),
        ("Budget Minimalist", "₹25 Lakh", "₹1.40 Cr", "4.6x", "Captures 68% of total risk reduction at only 25% of the capital budget.", c_amber),
        ("Max Reduction", "₹1.00 Cr", "₹2.70 Cr", "2.1x", "Exhaustive defense-in-depth minimizing residual financial exposure.", c_crimson)
    ]
    for i, (name, cost, red, rosi, note, col) in enumerate(portfolios):
        cx = Inches(0.8 + i * 2.95)
        add_card(s5, cx, Inches(2.2), Inches(2.8), Inches(4.5))
        pbox = s5.shapes.add_textbox(cx + Inches(0.2), Inches(2.4), Inches(2.4), Inches(4.1))
        ptf = pbox.text_frame
        ptf.word_wrap = True
        p = ptf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(10)
        
        metrics = [("Cost", cost), ("Risk Cut", red), ("ROSI", rosi)]
        for lbl, v in metrics:
            p = ptf.add_paragraph()
            p.text = f"{lbl}: {v}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = c_white
            p.space_after = Pt(4)
        
        p = ptf.add_paragraph()
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = c_muted
        p.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 6: Market Opportunity & Regulatory Tailwinds
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_bg(s6)
    add_header(s6, "Market Dynamics", "$18.5B TAM Driven by Mandatory Regulatory Catalysts", "Cyber risk quantification has transitioned from an elective exercise into a statutory imperative.")

    m_cards = [
        ("Total Addressable Market", "$18.5 Billion", "Global Cyber Risk Quantification & Security Analytics market by 2028 (CAGR 24.2%).", c_cyan),
        ("Serviceable Addressable Market", "$4.2 Billion", "BFSI, FinTech, Critical Infrastructure, and Healthcare across APAC, North America, and EMEA.", c_teal),
        ("Regulatory Driver: SEBI CSCRF", "Mandatory (2024–26)", "Mandates quantifiable cyber resilience limits, continuous telemetry, and third-party risk audits.", c_amber),
        ("Regulatory Driver: DPDP Act 2023", "Up to ₹250 Cr Fines", "Severe statutory penalties per data breach scenario in India, demanding precise risk quantification.", c_crimson)
    ]
    for i, (head, figure, desc, col) in enumerate(m_cards):
        col_idx = i % 2
        row_idx = i // 2
        cx = Inches(0.8 + col_idx * 5.95)
        cy = Inches(2.2 + row_idx * 2.4)
        add_card(s6, cx, cy, Inches(5.75), Inches(2.2))
        mbox = s6.shapes.add_textbox(cx + Inches(0.3), cy + Inches(0.2), Inches(5.15), Inches(1.8))
        mtf = mbox.text_frame
        mtf.word_wrap = True
        p1 = mtf.paragraphs[0]
        p1.text = head.upper()
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = mtf.add_paragraph()
        p2.text = figure
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = c_white
        p2.space_after = Pt(4)
        p3 = mtf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11.5)
        p3.font.color.rgb = c_muted

    # -------------------------------------------------------------------------
    # SLIDE 7: Competitive Advantage Matrix
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_bg(s7)
    add_header(s7, "Competitive Landscape", "Why CyberOptix Wins Against Incumbents", "Side-by-side comparison across quantification, causal rigor, and capital optimization.")

    # Table layout
    rows = 6
    cols = 5
    table_shape = s7.shapes.add_table(rows, cols, Inches(0.8), Inches(2.1), Inches(11.733), Inches(4.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.1)
    table.columns[3].width = Inches(2.1)
    table.columns[4].width = Inches(2.033)

    t_data = [
        ["Capability", "CyberOptix (CRIM-X)", "Safe Security (RiskLens)", "Kovrr", "Legacy GRC (Archer)"],
        ["Quantification Engine", "Continuous FAIR + MC", "FAIR / Monte Carlo", "Proprietary Statistical", "Subjective 5x5 Heatmaps"],
        ["Causal Machine Learning", "Causal DML (No Bias)", "Naive Correlation", "Industry Averages", "None"],
        ["Capital Optimizer", "5D Pareto Multi-Objective", "Simple Knapsack", "Scenario Comparison", "Static Task List"],
        ["Currency Denomination", "Native INR (Lakh/Cr)", "USD Only", "USD / EUR", "No Financial Values"],
        ["Tamper-Evident Proof", "SHA-256 Chained Hash", "No Cryptographic Proof", "None", "Modifiable Database"]
    ]
    for r_idx, r_data in enumerate(t_data):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(18, 29, 51)
            elif c_idx == 1:
                cell.fill.fore_color.rgb = RGBColor(12, 38, 66)
            else:
                cell.fill.fore_color.rgb = c_card
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if r_idx > 0 else 11.5)
                p.font.color.rgb = c_cyan if (c_idx == 1 and r_idx > 0) else c_white
                p.font.bold = (r_idx == 0 or c_idx <= 1)

    # -------------------------------------------------------------------------
    # SLIDE 8: Business Model & SaaS Unit Economics
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_bg(s8)
    add_header(s8, "Go-To-Market", "Predictable Enterprise SaaS Unit Economics", "Tiered annual contract values (ACV) delivering high retention and rapid payback.")

    tiers = [
        ("Enterprise Core", "₹18 Lakh / yr", "$22K ARR", "Continuous FAIR Engine, NIST/SEBI mapping, 10 telemetry connectors, PDF executive briefs.", c_cyan),
        ("Enterprise Apex", "₹38 Lakh / yr", "$46K ARR", "Full CRIM-X Causal DML, 5D Pareto Optimizer, Grounded AI Copilot, unlimited connectors.", c_teal),
        ("Global Financial Tier", "₹65 Lakh / yr", "$78K ARR", "Unlimited assets, on-premise/hybrid deployment, dedicated actuarial calibration, custom SLAs.", c_amber)
    ]
    for i, (name, acv_inr, acv_usd, desc, col) in enumerate(tiers):
        cx = Inches(0.8 + i * 3.95)
        add_card(s8, cx, Inches(2.2), Inches(3.75), Inches(4.5))
        tbox = s8.shapes.add_textbox(cx + Inches(0.3), Inches(2.4), Inches(3.15), Inches(4.0))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        p = ttf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(4)
        p = ttf.add_paragraph()
        p.text = acv_inr
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = c_white
        p = ttf.add_paragraph()
        p.text = f"({acv_usd})"
        p.font.size = Pt(12)
        p.font.color.rgb = c_muted
        p.space_after = Pt(12)
        p = ttf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = c_white
        p.space_after = Pt(16)
        p = ttf.add_paragraph()
        p.text = "Target Metrics:\n• 135%+ Net Revenue Retention\n• 82% Gross Margin\n• 7.2 Mo CAC Payback"
        p.font.size = Pt(11)
        p.font.color.rgb = c_muted

    # -------------------------------------------------------------------------
    # SLIDE 9: Traction & Architectural Validation
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_bg(s9)
    add_header(s9, "Execution Rigor", "Built & Validated to Institutional Standards", "Production-grade codebase with zero-regression automated test coverage.")

    proofs = [
        ("41/41 Tests Passing", "Complete automated pytest suite verifying FAIR engine, PuLP optimizer, CRIM-X causal models, and RBAC APIs.", c_teal),
        ("<0.25s Calculation Speed", "Vectorized NumPy Monte Carlo engine simulates 50,000 trials in real-time, enabling live what-if sliders.", c_cyan),
        ("74 kB Frontend Chunk", "Vite code-splitting with Gzip compression and lazy loading delivers instant sub-second page loads.", c_amber),
        ("Full Docker Stack", "Docker Compose, Nginx reverse proxy, async SQLite/PostgreSQL, and mobile PWA support.", c_crimson)
    ]
    for i, (stat, detail, col) in enumerate(proofs):
        col_idx = i % 2
        row_idx = i // 2
        cx = Inches(0.8 + col_idx * 5.95)
        cy = Inches(2.2 + row_idx * 2.4)
        add_card(s9, cx, cy, Inches(5.75), Inches(2.2))
        pbox = s9.shapes.add_textbox(cx + Inches(0.3), cy + Inches(0.2), Inches(5.15), Inches(1.8))
        ptf = pbox.text_frame
        ptf.word_wrap = True
        p1 = ptf.paragraphs[0]
        p1.text = stat
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.space_after = Pt(6)
        p2 = ptf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(12)
        p2.font.color.rgb = c_white

    # -------------------------------------------------------------------------
    # SLIDE 10: The Ask & Leadership Vision
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_bg(s10)
    add_header(s10, "Investment Ask", "Join Us in Defining the Future of Cyber Risk Capital", "Raising $2.5M Seed Extension to dominate the APAC & global financial cyber risk quantification market.")

    # Left: Allocation Card
    add_card(s10, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    abox = s10.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(4.0))
    atf = abox.text_frame
    atf.word_wrap = True
    p = atf.paragraphs[0]
    p.text = "$2.5M Capital Allocation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.space_after = Pt(12)
    allocs = [
        ("45% Engineering & Causal AI", "Expand consortium peer telemetry to 25,000 cross-industry trials."),
        ("35% Direct Sales & Enterprise GTM", "Scale enterprise direct sales across BFSI and critical infrastructure in India & APAC."),
        ("20% Regulatory Accreditation", "Formal SEBI, RBI, and NIST certification partnerships.")
    ]
    for h, d in allocs:
        p = atf.add_paragraph()
        p.text = f"• {h}: {d}"
        p.font.size = Pt(12)
        p.font.color.rgb = c_white
        p.space_after = Pt(8)

    # Right: Contact & Action Card
    add_card(s10, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), border_col=c_teal)
    cbox = s10.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(4.0))
    ctf = cbox.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    p.text = "Contact & Live Platform"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_teal
    p.space_after = Pt(14)
    info = [
        ("Live Enterprise Platform", "https://cyberoptix.onrender.com"),
        ("Interactive Swagger Docs", "https://cyberoptix.onrender.com/docs"),
        ("Executive Contact", "partners@cyberoptix.ai"),
        ("Corporate Headquarters", "Bengaluru & Mumbai, India"),
        ("Regulatory Grounding", "SEBI CSCRF • NIST CSF 2.0 • FAIR™")
    ]
    for lbl, val in info:
        p = ctf.add_paragraph()
        p.text = f"{lbl}: "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = c_muted
        p = ctf.add_paragraph()
        p.text = val
        p.font.size = Pt(12.5)
        p.font.color.rgb = c_white
        p.space_after = Pt(6)

    # Save presentation
    prs.save(output_filename)
    print(f"[+] Successfully generated widescreen presentation: {output_filename}")

if __name__ == "__main__":
    create_pitch_deck()
